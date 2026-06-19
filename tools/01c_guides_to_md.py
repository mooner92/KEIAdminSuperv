#!/usr/bin/env python3
"""
01c_guides_to_md.py — 연구행정 가이드(HWP/HWPX/PDF/PPTX) → Markdown (볼트 10_업무가이드/ 적재)

규정집(01)과 달리 이 묶음은 PDF·PPTX가 섞여 있고, 조문(제N조) 구조가 아니라 산문/표 위주다.
- HWP/HWPX: hwp-hwpx-parser (01과 동일, 표는 인라인 마크다운으로 추출)
- PDF: PyMuPDF(fitz) 텍스트 추출. 텍스트가 없으면(스캔 이미지) image-pdf로 표시 + 「TODO: 원문 확인」 플레이스홀더
- PPTX: python-pptx 슬라이드 텍스트/표 추출
- type: guide 프론트매터. 분류는 키워드로 규정집과 같은 버킷(3000_인사 등)에 매핑 → Explorer 필터 공유
- 파일명/슬러그는 볼트 전체와 충돌 안 나게 보장(규정 원문을 덮지 않음)

⛔ 절대 규칙: 원문 의역 금지. 변환 문구 보존, 표/그림 깨짐·오타만 사람이 검수. 항상 '검수상태: 미검수'.

테스트:  python 01c_guides_to_md.py --src research_rule_files --vault KEI-행정가이드 --dry-run
실행:    python 01c_guides_to_md.py --src research_rule_files --vault KEI-행정가이드
"""
import argparse
import datetime
import multiprocessing as mp
import queue as _queue
import re
import unicodedata
from pathlib import Path

FS_FORBIDDEN = re.compile(r'[\\/:*?"<>|\x00-\x1f]')

# 가이드 분류 키워드 → 규정집과 같은 버킷(미분류 가능). 위에서부터 첫 매치.
GUIDE_CATS = [
    ("4000_보수·여비", r"여비|출장|숙박|일비|보수|급여|연말정산|소득|수당"),
    ("3000_인사", r"인사|채용|복무|휴가|승진|전직|육아|근로시간|임금피크|평가|연수|교육|징계|갑질|윤리|진실성|행동강령|복리|후생|정년|퇴직|신입|길라잡이|당직|일용직|아르바이트|근태|출산|배우자"),
    ("6000_총무·보안·회계", r"예산|회계|법인카드|카드|경비|구매|계약|자산|보안|사이버|총무|경조사|편람|출판|포럼|업무추진비|물품|용역|재무|세금|계산서|영수증"),
    ("5000_연구·정보", r"연구|과제|R&?D|성과|논문|데이터|수탁|발주|정보화|시스템|환경포럼"),
    ("2000_감사·규정", r"감사|청렴|내부통제|부패|이해충돌"),
    ("1000_기관", r"기관|정관|조직|이사회"),
]
UNCLASSIFIED = "0000_미분류"


def nfc(s: str) -> str:
    return unicodedata.normalize("NFC", s)


def _valid(y, m, d=1):
    try:
        datetime.date(y, m, d)
        return True
    except ValueError:
        return False


def parse_date(text: str):
    """파일명에서 개정/작성일 추출 → 'YYYY-MM-DD' 또는 'YYYY-MM'. (01과 동일 로직)"""
    t = text
    m = re.search(r"(\d{4})\s*년\s*도?\s*(\d{1,2})\s*월\s*(\d{1,2})\s*일", t)
    if m:
        y, mo, d = map(int, m.groups())
        if _valid(y, mo, d):
            return f"{y:04d}-{mo:02d}-{d:02d}"
    m = re.search(r"(?<!\d)(\d{4})\.\s*(\d{1,2})\.\s*(\d{1,2})\.?", t)
    if m:
        y, mo, d = map(int, m.groups())
        if _valid(y, mo, d):
            return f"{y:04d}-{mo:02d}-{d:02d}"
    m = re.search(r"(?<!\d)(\d{2})\.\s*(\d{1,2})\.\s*(\d{1,2})\.", t)
    if m:
        yy, mo, d = map(int, m.groups())
        if _valid(2000 + yy, mo, d):
            return f"{2000 + yy:04d}-{mo:02d}-{d:02d}"
    for m in re.finditer(r"(?<!\d)(\d{4})(\d{2})(\d{2})(?!\d)", t):
        y, mo, d = int(m.group(1)), int(m.group(2)), int(m.group(3))
        if 1990 <= y <= 2099 and _valid(y, mo, d):
            return f"{y:04d}-{mo:02d}-{d:02d}"
    m = re.search(r"(?<!\d)(\d{4})\.\s*(\d{1,2})\.(?!\s*\d)", t)
    if m:
        y, mo = int(m.group(1)), int(m.group(2))
        if 1990 <= y <= 2099 and 1 <= mo <= 12:
            return f"{y:04d}-{mo:02d}"
    for m in re.finditer(r"(?<!\d)(\d{6})(?!\d)", t):
        s = m.group(1)
        yy, mo, d = int(s[:2]), int(s[2:4]), int(s[4:6])
        if _valid(2000 + yy, mo, d):
            return f"{2000 + yy:04d}-{mo:02d}-{d:02d}"
        y4, mo2 = int(s[:4]), int(s[4:6])
        if 1990 <= y4 <= 2099 and 1 <= mo2 <= 12:
            return f"{y4:04d}-{mo2:02d}"
    return None


def clean_title(stem: str) -> str:
    """파일명 → 사람이 읽을 제목(날짜·버전·장식 제거). (영문) 등 판본 표시는 유지."""
    t = stem
    t = re.sub(r"^[★☆■◆●▶]+\s*", "", t)
    t = re.sub(r"^\d{1,3}\.\s*", "", t)              # 리스트 마커 '50.'
    t = re.sub(r"\(?\s*\d{4}\s*년\s*도?\s*\d{1,2}\s*월\s*\d{1,2}\s*일\s*개?정?\s*\)?", " ", t)
    t = re.sub(r"\(?\s*\d{2,4}\.\s*\d{1,2}\.\s*\d{1,2}\.?\s*\)?", " ", t)
    t = re.sub(r"\(?\s*\d{4}\.\s*\d{1,2}\.?\s*\)?", " ", t)
    t = re.sub(r"[_(]?\s*\d{6,8}\.?\s*[_)]?", " ", t)
    t = re.sub(r"\(\s*\d+\s*\)", " ", t)             # 복사본 표시 (1)
    t = re.sub(r"\((?:최종|안|제정|개정|수정|개정판|[^)]*수정)\)", " ", t)
    t = t.replace("_", " ")
    t = re.sub(r"\s+", " ", t).strip(" -_·.")
    t = re.sub(r"\s*(전문개정|개정본|최종|제정|수정|개정판)$", "", t).strip(" -_·.")
    t = FS_FORBIDDEN.sub("", t)
    return t[:80] or stem[:80]


def classify(title: str, body: str) -> str:
    # 제목만으로 분류(본문 키워드는 노이즈가 커 3000_인사로 과집중됨). 못 맞추면 미분류.
    for cat, pat in GUIDE_CATS:
        if re.search(pat, title):
            return cat
    return UNCLASSIFIED


# ── HWP/HWPX (01과 동일: 별도 프로세스 + 하드 타임아웃) ──
def _hwp_worker(path_str, q):
    try:
        from hwp_hwpx_parser import Reader
        with Reader(path_str) as r:
            enc = r.is_encrypted
            enc = enc() if callable(enc) else enc
            if enc:
                q.put(("", "encrypted")); return
            body = (r.extract_text() or "").strip()
        q.put((body, "ok" if body else "empty"))
    except Exception as e:
        q.put((f"{type(e).__name__}: {e}", "error"))


def extract_hwp(path: Path, timeout: int):
    ctx = mp.get_context("fork")
    q = ctx.Queue()
    p = ctx.Process(target=_hwp_worker, args=(str(path), q), daemon=True)
    p.start()
    try:
        result = q.get(timeout=timeout)
    except _queue.Empty:
        result = ("", "timeout")
    except Exception as e:
        result = (f"{type(e).__name__}: {e}", "error")
    finally:
        if p.is_alive():
            p.terminate(); p.join(5)
            if p.is_alive():
                p.kill()
        else:
            p.join(1)
    return result


def extract_pdf(path: Path):
    try:
        import fitz
        parts = []
        with fitz.open(str(path)) as doc:
            for page in doc:
                parts.append(page.get_text("text"))
        body = re.sub(r"\n{3,}", "\n\n", "\n\n".join(parts)).strip()
        return (body, "ok" if len(body) >= 40 else "image-pdf")
    except Exception as e:
        return (f"{type(e).__name__}: {e}", "error")


def extract_pptx(path: Path):
    try:
        from pptx import Presentation
        parts = []
        prs = Presentation(str(path))
        for i, slide in enumerate(prs.slides, 1):
            slide_txt = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    slide_txt.append(shape.text_frame.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        slide_txt.append(" | ".join(c.text for c in row.cells))
            if slide_txt:
                parts.append(f"## 슬라이드 {i}\n\n" + "\n\n".join(slide_txt))
        body = "\n\n".join(parts).strip()
        return (body, "ok" if body else "empty")
    except Exception as e:
        return (f"{type(e).__name__}: {e}", "error")


def extract(path: Path, timeout: int):
    ext = path.suffix.lower()
    if ext in (".hwp", ".hwpx"):
        return extract_hwp(path, timeout)
    if ext == ".pdf":
        return extract_pdf(path)
    if ext == ".pptx":
        return extract_pptx(path)
    return ("", "skip")


def build_note(title, cat, date, original, body, image_pdf=False) -> str:
    if image_pdf or not body:
        body = "「TODO: 원문 확인」 — 스캔/이미지 PDF로 텍스트 자동추출 불가. 원본을 보고 본문을 채워야 합니다."
    fm = [
        "---",
        "type: guide",
        f'제목: "{title}"',
        f'분류: "{cat}"',
        '대상: "전직원"',
        "관련규정: []",
        "관련서식: []",
        f"개정일: {date}" if date else "개정일:",
        "최종검토일:",
        "검토자:",
        f'원본파일: "{original}"',
        "태그: []",
        "검수상태: 미검수",
        "---",
        "",
        f"# {title}",
        "",
        "> [!warning] 자동 변환(HWP/PDF/PPTX) — 의역 금지. 표/그림 깨짐과 오타만 검수 후 `검수상태: 검수완료`로.",
        "",
        body,
        "",
    ]
    return "\n".join(fm)


def main():
    ap = argparse.ArgumentParser(description="연구행정 가이드(HWP/HWPX/PDF/PPTX) → Markdown (10_업무가이드/)")
    ap.add_argument("--src", required=True)
    ap.add_argument("--vault", required=True)
    ap.add_argument("--dry-run", action="store_true")
    ap.add_argument("--limit", type=int, default=0)
    ap.add_argument("--timeout", type=int, default=120)
    args = ap.parse_args()

    vault = Path(args.vault)
    out_root = vault / "10_업무가이드"
    src = Path(args.src)
    files = sorted(
        list(src.glob("*.hwp")) + list(src.glob("*.hwpx"))
        + list(src.glob("*.pdf")) + list(src.glob("*.pptx"))
    )
    if args.limit:
        files = files[: args.limit]

    # 볼트 전체의 기존 파일 stem(규정 원문 포함) → 슬러그 충돌 방지(규정을 덮지 않음)
    existing = {md.stem for md in vault.rglob("*.md")}
    print(f"{len(files)}개 파일 {'미리보기' if args.dry_run else '변환'} 시작 · 기존 볼트 문서 {len(existing)}개")

    stats = {"ok": 0, "image_pdf": 0, "skip": 0}
    by_cat: dict[str, int] = {}
    by_ext: dict[str, int] = {}
    rows = []

    for i, f in enumerate(files, 1):
        ext = f.suffix.lower()
        stem = nfc(f.stem)
        body, status = extract(f, args.timeout)
        image_pdf = status == "image-pdf"
        ok = status in ("ok", "image-pdf")
        print(f"  [{i:>3}/{len(files)}] {status:<10} {ext:<6} {f.name}", flush=True)
        if not ok:
            stats["skip"] += 1
            continue

        title = clean_title(stem)
        cat = classify(title, "")
        date = parse_date(stem)

        # 슬러그 유일성: 기존 볼트 + 이번 배치와 충돌 시 _2, _3 …
        base = title
        slug = base
        n = 1
        while slug in existing:
            n += 1
            slug = f"{base}_{n}"
        existing.add(slug)

        dest = out_root / cat / f"{slug}.md"
        by_cat[cat] = by_cat.get(cat, 0) + 1
        by_ext[ext] = by_ext.get(ext, 0) + 1
        rows.append((ext, "image-pdf" if image_pdf else "ok", date or "—", cat, title[:40]))
        stats["image_pdf" if image_pdf else "ok"] += 1

        if not args.dry_run:
            dest.parent.mkdir(parents=True, exist_ok=True)
            dest.write_text(build_note(title, cat, date, f.name, body, image_pdf), encoding="utf-8")

    print("\n" + "=" * 96)
    print(f"{'확장':<6}{'상태':<10}{'개정일':<12}{'분류':<20}제목")
    print("-" * 96)
    for ext, st, date, cat, title in rows:
        print(f"{ext:<6}{st:<10}{date:<12}{cat:<20}{title}")
    print("=" * 96)
    print(f"성공 {stats['ok']} · 이미지PDF(플레이스홀더) {stats['image_pdf']} · 건너뜀 {stats['skip']}")
    print("확장자별:", ", ".join(f"{k} {v}" for k, v in sorted(by_ext.items())))
    print("분류별:", ", ".join(f"{k} {v}" for k, v in sorted(by_cat.items())))
    if args.dry_run:
        print("\n(dry-run: 파일 미작성)")


if __name__ == "__main__":
    main()
